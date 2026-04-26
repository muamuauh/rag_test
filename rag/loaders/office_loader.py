from __future__ import annotations

from pathlib import Path

from docx import Document
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

from rag.config import settings
from rag.loaders.base import image_hash, is_significant_image
from rag.schema import Chunk


def load_docx(path: str | Path) -> list[Chunk]:
    """Word 文档：所有段落合并为一个文本 Chunk（后续按 chunker 二次切分）。
    内嵌图片提取为 image_caption 占位 Chunk。"""
    path = Path(path)
    chunks: list[Chunk] = []
    doc = Document(path)

    paragraphs = [p.text.strip() for p in doc.paragraphs if p.text.strip()]
    if paragraphs:
        chunks.append(
            Chunk(
                text="\n".join(paragraphs),
                source_path=str(path),
                chunk_type="text",
                metadata={"file_name": path.name},
            )
        )

    images_dir = settings.data_images_dir
    images_dir.mkdir(parents=True, exist_ok=True)
    for rel in doc.part.rels.values():
        if "image" not in rel.reltype:
            continue
        try:
            data: bytes = rel.target_part.blob
            ext = Path(rel.target_part.partname).suffix.lstrip(".") or "png"
        except Exception:
            continue
        if not is_significant_image(data):
            continue
        h = image_hash(data)
        img_path = images_dir / f"{h}.{ext}"
        if not img_path.exists():
            img_path.write_bytes(data)
        chunks.append(
            Chunk(
                text="",
                source_path=str(path),
                chunk_type="image_caption",
                image_path=str(img_path),
                metadata={"file_name": path.name, "image_hash": h},
            )
        )

    return chunks


def load_pptx(path: str | Path) -> list[Chunk]:
    """PPT：每张幻灯片一个文本 Chunk + 幻灯片中图片各一个 image_caption Chunk。"""
    path = Path(path)
    chunks: list[Chunk] = []
    prs = Presentation(path)
    images_dir = settings.data_images_dir
    images_dir.mkdir(parents=True, exist_ok=True)

    for slide_idx, slide in enumerate(prs.slides):
        slide_text_parts: list[str] = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                txt = shape.text_frame.text.strip()
                if txt:
                    slide_text_parts.append(txt)
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                try:
                    image = shape.image
                    data: bytes = image.blob
                    ext: str = image.ext or "png"
                except Exception:
                    continue
                if not is_significant_image(data):
                    continue
                h = image_hash(data)
                img_path = images_dir / f"{h}.{ext}"
                if not img_path.exists():
                    img_path.write_bytes(data)
                chunks.append(
                    Chunk(
                        text="",
                        source_path=str(path),
                        chunk_type="image_caption",
                        page=slide_idx + 1,
                        image_path=str(img_path),
                        metadata={"file_name": path.name, "image_hash": h},
                    )
                )
        if slide_text_parts:
            chunks.append(
                Chunk(
                    text="\n".join(slide_text_parts),
                    source_path=str(path),
                    chunk_type="text",
                    page=slide_idx + 1,
                    metadata={"file_name": path.name},
                )
            )

    return chunks
